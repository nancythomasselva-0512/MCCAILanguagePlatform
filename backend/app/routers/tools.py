from fastapi import APIRouter, Depends, HTTPException, status, UploadFile, File, Form
from sqlalchemy.orm import Session
from app.core.database import get_db
from app.dependencies.auth import get_current_user, get_current_tenant_context, check_tenant_access
from app.models.models import User, Tenant, UsageTracking, ProviderConfiguration, TranscriptionHistory, TranslationHistory, TtsHistory, SubscriptionPlan, FeatureProviderMapping
from app.schemas.schemas import TranscriptionResponse, TranslationResponse, TtsResponse
from app.core.security import decrypt_data
from app.core.config import settings
from app.ai.provider_manager import provider_manager
import datetime
import requests
import tempfile
import os
import base64
from typing import List

router = APIRouter(prefix="/tools", tags=["AI Processing Tools"], dependencies=[Depends(check_tenant_access)])

# Endpoint to fetch the currently active provider for a given feature (used by frontend badge)
@router.get("/active-providers")
def get_active_providers(db: Session = Depends(get_db)):
    # Returns a map of feature_name -> active provider_name
    mappings = db.query(FeatureProviderMapping).filter(FeatureProviderMapping.is_enabled == True).order_by(FeatureProviderMapping.priority.asc()).all()
    
    result = {}
    for m in mappings:
        if m.feature_name not in result:
            result[m.feature_name] = m.provider_name
    return result

def get_global_provider_for_feature(db: Session, tenant: Tenant, feature_name: str) -> tuple[str, str]:
    mappings = db.query(FeatureProviderMapping).filter(
        FeatureProviderMapping.feature_name == feature_name,
        FeatureProviderMapping.is_enabled == True
    ).order_by(FeatureProviderMapping.priority.asc()).all()
    
    for mapping in mappings:
        api_key = resolve_api_key(db, tenant, mapping.provider_name)
        # Assuming if there's an api_key or if it's a provider that doesn't need an API key
        return mapping.provider_name, api_key
        
    return "openai", resolve_api_key(db, tenant, "openai") # Safe fallback

# Helper to check usage limit
def verify_limit(db: Session, tenant: Tenant, metric: str, incremental_amount: float):
    if not tenant:
        return  # Bypass if no tenant (e.g. Super Admin without context)
        
    plan = db.query(SubscriptionPlan).filter(SubscriptionPlan.id == tenant.plan_id).first()
    usage = db.query(UsageTracking).filter(UsageTracking.tenant_id == tenant.id).first()
    
    if not plan or not usage:
        return  # Bypass if no plan details are present
        
    if metric == "transcription":
        if usage.audio_minutes_used + incremental_amount > plan.transcription_limit:
            raise HTTPException(
                status_code=status.HTTP_402_PAYMENT_REQUIRED,
                detail=f"Subscription limit exceeded. Transcription limit is {plan.transcription_limit} mins, you used {round(usage.audio_minutes_used, 2)} mins."
            )
    elif metric == "translation":
        if usage.translation_chars_used + incremental_amount > plan.translation_limit:
            raise HTTPException(
                status_code=status.HTTP_402_PAYMENT_REQUIRED,
                detail=f"Subscription limit exceeded. Translation limit is {plan.translation_limit} chars, you used {usage.translation_chars_used} chars."
            )
    elif metric == "tts":
        if usage.tts_chars_used + incremental_amount > plan.tts_limit:
            raise HTTPException(
                status_code=status.HTTP_402_PAYMENT_REQUIRED,
                detail=f"Subscription limit exceeded. TTS limit is {plan.tts_limit} chars, you used {usage.tts_chars_used} chars."
            )

# Helper to resolve API Key: checks Tenant settings first, falls back to Global Settings
def resolve_api_key(db: Session, tenant: Tenant, provider_name: str) -> str:
    # 1. Check Tenant config
    if tenant:
        tenant_config = db.query(ProviderConfiguration).filter(
            ProviderConfiguration.tenant_id == tenant.id,
            ProviderConfiguration.provider_name == provider_name,
            ProviderConfiguration.is_enabled == True
        ).first()
        if tenant_config and tenant_config.credentials_encrypted:
            return decrypt_data(tenant_config.credentials_encrypted)
        
    # 2. Fall back to Global config
    global_config = db.query(ProviderConfiguration).filter(
        ProviderConfiguration.tenant_id == None,
        ProviderConfiguration.provider_name == provider_name,
        ProviderConfiguration.is_enabled == True
    ).first()
    if global_config and global_config.credentials_encrypted:
        return decrypt_data(global_config.credentials_encrypted)
        
    # 3. Else fallback to environment key variable if present
    env_var_map = {
        "openai": "OPENAI_API_KEY",
        "deepgram": "DEEPGRAM_API_KEY",
        "elevenlabs": "ELEVENLABS_API_KEY"
    }
    env_key = env_var_map.get(provider_name)
    if env_key:
        val = os.getenv(env_key)
        if val:
            return val
            
    return ""

# Increment usage helper
def increment_usage(db: Session, tenant_id: str, metric: str, amount: float):
    if not tenant_id:
        return
    usage = db.query(UsageTracking).filter(UsageTracking.tenant_id == tenant_id).first()
    if not usage:
        usage = UsageTracking(tenant_id=tenant_id)
        db.add(usage)
    
    if metric == "transcription":
        usage.audio_minutes_used += amount
    elif metric == "translation":
        usage.translation_chars_used += int(amount)
    elif metric == "tts":
        usage.tts_chars_used += int(amount)
        
    usage.api_calls_used += 1
    db.commit()

# --- TOOL 1: TEXT TRANSLATION ---
@router.post("/translate")
def translate_text(
    text: str = Form(...),
    source_lang: str = Form("Auto Detect"),
    target_lang: str = Form("English"),
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
    tenant: Tenant = Depends(get_current_tenant_context)
):
    char_len = len(text)
    verify_limit(db, tenant, "translation", char_len)

    # Build translation prompt
    if source_lang and source_lang != "Auto Detect":
        system_prompt = (
            f"You are a professional translator. "
            f"Translate the following text from {source_lang} to {target_lang}. "
            f"Output ONLY the translated text with no quotes, explanations, or extra commentary."
        )
    else:
        system_prompt = (
            f"You are a professional translator. "
            f"Translate the following text to {target_lang}. "
            f"Output ONLY the translated text with no quotes, explanations, or extra commentary."
        )

    # Delegate to provider manager — automatic fallback + circuit breaker
    translated_text = provider_manager.call_llm(
        system_prompt=system_prompt,
        user_message=text,
        feature="translation",
        db=db,
        tenant=tenant,
    )

    # Persist history and usage
    history = TranslationHistory(
        tenant_id=tenant.id if tenant else None,
        user_id=user.id,
        source_text=text,
        translated_text=translated_text,
        source_lang=source_lang if source_lang != "Auto Detect" else "Auto",
        target_lang=target_lang,
        provider="auto",  # actual provider tracked in ProviderLog
    )
    db.add(history)
    db.commit()
    increment_usage(db, tenant.id if tenant else None, "translation", char_len)

    return {
        "text": translated_text,
        "source_lang": source_lang,
        "target_lang": target_lang,
        "detected_lang": "en",
    }

# --- TOOL 2: VOICE & AUDIO TRANSCRIPTION ---
@router.post("/transcribe")
async def transcribe_audio(
    file: UploadFile = File(...),
    model: str = Form("base"),
    language: str = Form("en"),
    feature_name: str = Form("Audio To Text"),
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
    tenant: Tenant = Depends(get_current_tenant_context)
):
    audio_bytes = await file.read()
    file_size = len(audio_bytes)
    estimated_minutes = max(0.1, (file_size / (1024 * 1024)) * 0.5)
    verify_limit(db, tenant, "transcription", estimated_minutes)

    # Delegate to provider manager — tries local-whisper → deepgram → openai-stt
    result = provider_manager.call_stt(
        audio_bytes=audio_bytes,
        filename=file.filename or "audio.wav",
        language=language,
        model=model,
        feature=feature_name,
        db=db,
        tenant=tenant,
    )
    transcript_text = result.get("text", "")
    segments_data   = result.get("segments", [])

    history = TranscriptionHistory(
        tenant_id=tenant.id if tenant else None,
        user_id=user.id,
        file_name=file.filename or "voice-recording.wav",
        file_size=file_size,
        duration_seconds=estimated_minutes * 60,
        transcript_text=transcript_text,
        provider="auto",
    )
    db.add(history)
    db.commit()
    increment_usage(db, tenant.id if tenant else None, "transcription", estimated_minutes)

    if not segments_data:
        segments_data = [{"timestamp": "00:00", "text": transcript_text, "start": 0, "end": estimated_minutes * 60}]

    return {"text": transcript_text, "provider": "auto", "segments": segments_data}


# --- TOOL 3: TEXT-TO-SPEECH SYNTHESIS ---
@router.post("/synthesize")
def synthesize_speech(
    text: str = Form(...),
    voice: str = Form("alloy"),
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
    tenant: Tenant = Depends(get_current_tenant_context)
):
    char_len = len(text)
    verify_limit(db, tenant, "tts", char_len)

    # Delegate to provider manager — tries elevenlabs → openai-tts
    audio_bytes = provider_manager.call_tts(
        text=text,
        voice=voice,
        feature="Text To Speech",
        db=db,
        tenant=tenant,
    )
    audio_url = f"data:audio/mpeg;base64,{base64.b64encode(audio_bytes).decode('utf-8')}"

    history = TtsHistory(
        tenant_id=tenant.id if tenant else None,
        user_id=user.id,
        text=text,
        voice_name=voice,
        characters_count=char_len,
        file_path="",
        provider="auto",
    )
    db.add(history)
    db.commit()
    increment_usage(db, tenant.id if tenant else None, "tts", char_len)

    return {
        "status": "success",
        "characters": char_len,
        "voice": voice,
        "provider": "auto",
        "audio_url": audio_url,
    }

# --- USER TRANSLATION HISTORY LOG ---
@router.get("/history/translations", response_model=List[TranslationResponse])
def get_translations_history(
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
    tenant: Tenant = Depends(get_current_tenant_context)
):
    # Enforce strict multi-tenant filtering (only fetch logs within this tenant)
    if user.role == "super_admin":
        return db.query(TranslationHistory).order_by(TranslationHistory.created_at.desc()).all()
    elif tenant:
        return db.query(TranslationHistory).filter(TranslationHistory.tenant_id == tenant.id).order_by(TranslationHistory.created_at.desc()).all()
    else:
        return db.query(TranslationHistory).filter(TranslationHistory.user_id == user.id, TranslationHistory.tenant_id == None).order_by(TranslationHistory.created_at.desc()).all()

# --- USER TRANSCRIPTION HISTORY LOG ---
@router.get("/history/transcriptions", response_model=List[TranscriptionResponse])
def get_transcriptions_history(
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
    tenant: Tenant = Depends(get_current_tenant_context)
):
    if user.role == "super_admin":
        return db.query(TranscriptionHistory).order_by(TranscriptionHistory.created_at.desc()).all()
    elif tenant:
        return db.query(TranscriptionHistory).filter(TranscriptionHistory.tenant_id == tenant.id).order_by(TranscriptionHistory.created_at.desc()).all()
    else:
        return db.query(TranscriptionHistory).filter(TranscriptionHistory.user_id == user.id, TranscriptionHistory.tenant_id == None).order_by(TranscriptionHistory.created_at.desc()).all()

# --- USER TTS HISTORY LOG ---
@router.get("/history/tts", response_model=List[TtsResponse])
def get_tts_history(
    db: Session = Depends(get_db),
    user: User = Depends(get_current_user),
    tenant: Tenant = Depends(get_current_tenant_context)
):
    if user.role == "super_admin":
        return db.query(TtsHistory).order_by(TtsHistory.created_at.desc()).all()
    elif tenant:
        return db.query(TtsHistory).filter(TtsHistory.tenant_id == tenant.id).order_by(TtsHistory.created_at.desc()).all()
    else:
        return db.query(TtsHistory).filter(TtsHistory.user_id == user.id, TtsHistory.tenant_id == None).order_by(TtsHistory.created_at.desc()).all()

# --- TOOL 4: TEXT EXTRACTION FROM FILE ---
@router.post("/extract-text")
async def extract_text_from_upload(
    file: UploadFile = File(...),
    user: User = Depends(get_current_user)
):
    try:
        content = await file.read()
        filename = file.filename.lower()
        ext = filename.split('.')[-1]
        text = ""

        if ext in ["txt", "csv", "json"]:
            text = content.decode("utf-8", errors="ignore")
        elif ext == "pdf":
            import PyPDF2
            import io
            reader = PyPDF2.PdfReader(io.BytesIO(content))
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n\n"
        elif ext in ["doc", "docx"]:
            import docx
            import io
            doc = docx.Document(io.BytesIO(content))
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n\n"
        elif ext in ["xls", "xlsx"]:
            import pandas as pd
            import io
            df = pd.read_excel(io.BytesIO(content))
            text = df.to_string(index=False)
        elif ext in ["ppt", "pptx"]:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(content))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n\n"
        elif ext in ["png", "jpg", "jpeg", "webp"]:
            import pytesseract
            from PIL import Image
            import io
            image = Image.open(io.BytesIO(content))
            text = pytesseract.image_to_string(image)
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")

        if not text.strip():
            raise HTTPException(status_code=400, detail="No readable text found in the file.")

        return {"text": text.strip()}

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting text: {str(e)}")
